"""Procesador de archivos para análisis con IA.

Extrae contenido de diferentes tipos de archivos (Word, PDF, Excel, etc.)
y lo prepara para envío a servicios de IA.

Author: Gemini AI Assistant
Date: 31 octubre 2025
"""

import io
import logging
from pathlib import Path
from typing import Any, BinaryIO


# Imports condicionales (se validan en runtime)
try:
    import PyPDF2

    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl

    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from src.services.ai.exceptions import FileProcessingError


logger = logging.getLogger(__name__)


class FileProcessor:
    """Procesador universal de archivos para análisis con IA.

    Soporta:
    - Documentos: PDF, DOCX, TXT, MD
    - Hojas de cálculo: XLSX, CSV
    - Presentaciones: PPTX
    - Código: Python, JavaScript, Java, etc.
    - Imágenes: JPG, PNG (vía Gemini Vision API)
    """

    SUPPORTED_TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".py",
        ".js",
        ".java",
        ".cpp",
        ".c",
        ".h",
        ".css",
        ".html",
        ".xml",
        ".json",
        ".yaml",
        ".yml",
        ".sh",
        ".sql",
        ".r",
        ".m",
        ".swift",
        ".go",
        ".rs",
        ".kt",
        ".ts",
    }

    SUPPORTED_MIME_TYPES = {
        "text/plain": "text",
        "text/markdown": "text",
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "text/csv": "csv",
        "image/jpeg": "image",
        "image/png": "image",
        "image/webp": "image",
        "application/x-python-code": "text",
        "text/x-python": "text",
        "text/javascript": "text",
        "application/json": "text",
    }

    @staticmethod
    def detectar_tipo_archivo(nombre_archivo: str, mime_type: str | None = None) -> str:
        """Detecta el tipo de archivo basándose en extensión y MIME type.

        Args:
            nombre_archivo: Nombre del archivo
            mime_type: MIME type del archivo

        Returns:
            str: Tipo detectado ('pdf', 'docx', 'text', etc.)

        Raises:
            FileProcessingError: Si el tipo no es soportado
        """
        # Intentar por MIME type primero
        if mime_type and mime_type in FileProcessor.SUPPORTED_MIME_TYPES:
            return FileProcessor.SUPPORTED_MIME_TYPES[mime_type]

        # Intentar por extensión
        extension = Path(nombre_archivo).suffix.lower()

        if extension in FileProcessor.SUPPORTED_TEXT_EXTENSIONS:
            return "text"
        if extension == ".pdf":
            return "pdf"
        if extension == ".docx":
            return "docx"
        if extension == ".xlsx":
            return "xlsx"
        if extension == ".pptx":
            return "pptx"
        if extension == ".csv":
            return "csv"
        if extension in {".jpg", ".jpeg", ".png", ".webp"}:
            return "image"

        msg = f"Tipo de archivo no soportado: {extension}"
        raise FileProcessingError(
            msg,
            file_name=nombre_archivo,
            file_type=mime_type or extension,
        )

    @staticmethod
    def extraer_contenido(
        archivo: BinaryIO, nombre_archivo: str, mime_type: str | None = None
    ) -> str:
        """Extrae el contenido de texto de un archivo.

        Args:
            archivo: Stream binario del archivo
            nombre_archivo: Nombre del archivo
            mime_type: MIME type del archivo

        Returns:
            str: Contenido extraído en texto plano

        Raises:
            FileProcessingError: Si no se puede procesar el archivo
        """
        tipo = FileProcessor.detectar_tipo_archivo(nombre_archivo, mime_type)

        try:
            if tipo == "text":
                return FileProcessor._extraer_texto_plano(archivo)
            if tipo == "pdf":
                return FileProcessor._extraer_pdf(archivo)
            if tipo == "docx":
                return FileProcessor._extraer_docx(archivo)
            if tipo == "xlsx":
                return FileProcessor._extraer_xlsx(archivo)
            if tipo == "pptx":
                return FileProcessor._extraer_pptx(archivo)
            if tipo == "csv":
                return FileProcessor._extraer_csv(archivo)
            if tipo == "image":
                return FileProcessor._extraer_imagen(archivo, nombre_archivo)
            msg = f"Tipo de archivo '{tipo}' no implementado"
            raise FileProcessingError(msg, file_name=nombre_archivo)

        except FileProcessingError:
            raise  # Re-lanzar excepciones propias
        except Exception as e:
            logger.exception(f"Error procesando {nombre_archivo}: {e!s}")
            msg = f"Error al extraer contenido: {e!s}"
            raise FileProcessingError(
                msg,
                file_name=nombre_archivo,
                file_type=tipo,
            ) from e

    @staticmethod
    def _extraer_texto_plano(archivo: BinaryIO) -> str:
        """Extrae contenido de archivos de texto plano."""
        contenido = archivo.read()

        # Intentar diferentes encodings
        for encoding in ["utf-8", "latin-1", "cp1252", "iso-8859-1"]:
            try:
                return contenido.decode(encoding)
            except UnicodeDecodeError:
                continue

        msg = "No se pudo decodificar el archivo de texto"
        raise FileProcessingError(msg)

    @staticmethod
    def _extraer_pdf(archivo: BinaryIO) -> str:
        """Extrae texto de archivos PDF."""
        if not HAS_PDF:
            msg = "PyPDF2 no instalado. Instala con: pip install PyPDF2"
            raise FileProcessingError(msg)

        try:
            pdf_reader = PyPDF2.PdfReader(archivo)
            texto = []

            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    texto.append(f"--- Página {page_num} ---\n{page_text}")

            if not texto:
                msg = "No se pudo extraer texto del PDF"
                raise FileProcessingError(msg)

            return "\n\n".join(texto)

        except Exception as e:
            msg = f"Error procesando PDF: {e!s}"
            raise FileProcessingError(msg) from e

    @staticmethod
    def _extraer_docx(archivo: BinaryIO) -> str:
        """Extrae texto de archivos Word (.docx)."""
        if not HAS_DOCX:
            msg = "python-docx no instalado. Instala con: pip install python-docx"
            raise FileProcessingError(msg)

        try:
            doc = Document(archivo)
            texto = []

            # Extraer párrafos
            for para in doc.paragraphs:
                if para.text.strip():
                    texto.append(para.text)

            # Extraer tablas
            for table in doc.tables:
                tabla_texto = []
                for row in table.rows:
                    row_texto = " | ".join(cell.text for cell in row.cells)
                    tabla_texto.append(row_texto)
                if tabla_texto:
                    texto.append("\n[TABLA]\n" + "\n".join(tabla_texto) + "\n")

            if not texto:
                msg = "El documento Word está vacío"
                raise FileProcessingError(msg)

            return "\n\n".join(texto)

        except Exception as e:
            msg = f"Error procesando DOCX: {e!s}"
            raise FileProcessingError(msg) from e

    @staticmethod
    def _extraer_xlsx(archivo: BinaryIO) -> str:
        """Extrae datos de archivos Excel (.xlsx)."""
        if not HAS_EXCEL:
            msg = "openpyxl no instalado. Instala con: pip install openpyxl"
            raise FileProcessingError(msg)

        try:
            wb = openpyxl.load_workbook(archivo, data_only=True)
            texto = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                texto.append(f"--- Hoja: {sheet_name} ---\n")

                for row in sheet.iter_rows(values_only=True):
                    # Filtrar valores None y convertir a strings
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        texto.append(" | ".join(row_values))

            if len(texto) <= 1:  # Solo header de hoja
                msg = "El archivo Excel está vacío"
                raise FileProcessingError(msg)

            return "\n".join(texto)

        except Exception as e:
            msg = f"Error procesando XLSX: {e!s}"
            raise FileProcessingError(msg) from e

    @staticmethod
    def _extraer_pptx(archivo: BinaryIO) -> str:
        """Extrae texto de archivos PowerPoint (.pptx)."""
        if not HAS_PPTX:
            msg = "python-pptx no instalado. Instala con: pip install python-pptx"
            raise FileProcessingError(msg)

        try:
            prs = Presentation(archivo)
            texto = []

            for slide_num, slide in enumerate(prs.slides, 1):
                texto.append(f"--- Diapositiva {slide_num} ---")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texto.append(shape.text)

                texto.append("")  # Línea en blanco entre slides

            if not texto:
                msg = "La presentación está vacía"
                raise FileProcessingError(msg)

            return "\n".join(texto)

        except Exception as e:
            msg = f"Error procesando PPTX: {e!s}"
            raise FileProcessingError(msg) from e

    @staticmethod
    def _extraer_csv(archivo: BinaryIO) -> str:
        """Extrae datos de archivos CSV."""
        import csv

        try:
            # Detectar encoding
            muestra = archivo.read(4096)
            archivo.seek(0)

            for encoding in ["utf-8", "latin-1", "cp1252"]:
                try:
                    muestra.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue

            # Leer CSV
            contenido = archivo.read().decode(encoding)
            csv_reader = csv.reader(io.StringIO(contenido))

            filas = []
            for row in csv_reader:
                if row:  # Ignorar filas vacías
                    filas.append(" | ".join(row))

            if not filas:
                msg = "El archivo CSV está vacío"
                raise FileProcessingError(msg)

            return "\n".join(filas)

        except Exception as e:
            msg = f"Error procesando CSV: {e!s}"
            raise FileProcessingError(msg) from e

    @staticmethod
    def _extraer_imagen(archivo: BinaryIO, nombre_archivo: str) -> str:
        """Placeholder para procesamiento de imágenes.

        Las imágenes se procesan directamente por Gemini Vision API,
        este método retorna metadata.
        """
        return (
            f"[IMAGEN: {nombre_archivo}]\n"
            "Este archivo de imagen será procesado directamente por Gemini Vision API.\n"
            "El análisis visual se incluirá en la retroalimentación."
        )

    @staticmethod
    def validar_archivo(
        archivo: BinaryIO, nombre_archivo: str, max_size_mb: int = 20
    ) -> dict[str, Any]:
        """Valida que un archivo cumpla con restricciones.

        Args:
            archivo: Stream del archivo
            nombre_archivo: Nombre del archivo
            max_size_mb: Tamaño máximo en MB

        Returns:
            Dict con resultado de validación:
            {
                "valido": bool,
                "razones": List[str],
                "tamano_mb": float,
                "tipo": str
            }
        """
        razones = []

        # Verificar tamaño
        archivo.seek(0, 2)  # Ir al final
        tamano_bytes = archivo.tell()
        archivo.seek(0)  # Volver al inicio

        tamano_mb = tamano_bytes / (1024 * 1024)

        if tamano_mb > max_size_mb:
            razones.append(
                f"Archivo muy grande: {tamano_mb:.2f}MB (máximo {max_size_mb}MB)"
            )

        # Verificar tipo
        try:
            tipo = FileProcessor.detectar_tipo_archivo(nombre_archivo)
        except FileProcessingError as e:
            razones.append(str(e))
            tipo = "desconocido"

        # Verificar que no esté vacío
        if tamano_bytes == 0:
            razones.append("El archivo está vacío")

        return {
            "valido": len(razones) == 0,
            "razones": razones,
            "tamano_mb": round(tamano_mb, 2),
            "tipo": tipo,
        }
